"""documents package.

Implementation lives in :mod:`raghub.helper` (documents); local entry-point modules: ['parser'].
"""

from __future__ import annotations

# --- parser.py content ---
# --- parser.py content ---
import io
from abc import ABC, abstractmethod
from dataclasses import dataclass
from importlib import import_module
from io import BytesIO
from pathlib import Path
from typing import Any

from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from PIL import Image as PillowImage
from pptx import Presentation
from pypdf import PdfReader

from raghub.helper.documents import (
    EQUATION_BLOCK_RE,
    FENCE_RE,
    HEADING_RE,
    IMAGE_RE,
    INLINE_EQUATION_RE,
    MAGIC_BYTES,
    MARKER_AVAILABLE,
    MIME_TYPES_BY_EXTENSION,
    TABLE_LINE_RE,
    ChunkingPlan,
    DocumentLifecycleManager,
    MarkdownSection,
    MarkerConverter,
    MarkerImportError,
    MarkerPdfConverter,
    PlainTextConverter,
    build_chunk_records,
    build_marker_converter,
    chunk_words,
    convert_path,
    datetime_now_utc,
    detect_mime_type,
    extract_pdf_metadata,
    extract_pdf_pages,
    extract_pdf_text,
    extract_text_from_content,
    looks_like_pdf,
    markdown_to_document_blocks,
    marker_create_model_dict,
    marker_text_from_rendered,
    new_version,
    normalise_markdown,
    normalize_text,
    select_converter_for_path,
    self_module,
    validate_upload,
)
from raghub.utils import capture


@dataclass(frozen=True)
class Section:
    """A single parsed chunk of a document.

    Attributes:
        section_index: 0-based ordinal of this section within the file
            (e.g. PDF page number minus 1, or 0 for whole-file formats).
        source_location: Human-readable location string used as the
            ``source_location`` field on chunk records.
        text: The extracted text content.
        metadata: Format-specific metadata (e.g. PDF ``/Title``,
            ``/Author``). Optional; defaults to an empty dict.
    """

    section_index: int
    source_location: str
    text: str
    metadata: dict[str, Any]


class File(ABC):
    """Abstract base for all document parsers."""

    @abstractmethod
    def parse(self, file_bytes: bytes, file_name: str, mime_type: str) -> list[Section]:
        """Parse ``file_bytes`` into a list of :class:`Section`.

        Args:
            file_bytes: Raw file contents.
            file_name: Original filename; useful for format-specific
                hints (e.g. extension-based fallbacks).
            mime_type: The MIME type reported by the validator.

        Returns:
            A list of :class:`Section` objects. Empty when the
            parser finds no extractable text.
        """


class Pdf(File):
    """PDF file using :mod:`pypdf`."""

    def parse(self, file_bytes: bytes, file_name: str, mime_type: str) -> list[Section]:
        """Parse a PDF into one section per page.

        Args:
            file_bytes: Raw PDF bytes.
            file_name: Original filename (unused beyond diagnostics).
            mime_type: MIME type (unused; the parser always uses
                :class:`pypdf.PdfReader`).

        Returns:
            A list of :class:`Section`, one per page. Empty
            strings are returned for image-only pages rather than
            raising. The section's ``section_index`` is the 1-based
            page number and ``source_location`` is ``"page N"``. The
            metadata dict contains ``width`` and ``height`` (from the
            page's media box) when available.
        """
        reader = PdfReader(BytesIO(file_bytes))
        sections: list[Section] = []
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            sections.append(
                Section(
                    section_index=i,
                    source_location=f"page {i}",
                    text=text,
                    metadata={
                        "width": page.mediabox.width,
                        "height": page.mediabox.height,
                    }
                    if page.mediabox
                    else {},
                )
            )
        return sections


class HTML(File):
    """HTML file using :mod:`BeautifulSoup`."""

    def parse(self, file_bytes: bytes, file_name: str, mime_type: str) -> list[Section]:
        """Parse an HTML document into a single section.

        Args:
            file_bytes: Raw HTML bytes.
            file_name: Original filename (unused).
            mime_type: MIME type (unused).

        Returns:
            A single-element list with one :class:`Section`
            containing the concatenated body text. ``section_index``
            is 0 and ``source_location`` is ``"full file"``. The
            section metadata includes a ``headings`` list with the
            text of every ``<h1>``, ``<h2>``, and ``<h3>`` element.
        """
        soup = BeautifulSoup(file_bytes, "lxml")
        body = soup.find("body") or soup
        text = body.get_text(separator=" ", strip=True)
        return [
            Section(
                section_index=0,
                source_location="full file",
                text=text,
                metadata={
                    "headings": [h.get_text(strip=True) for h in soup.find_all(["h1", "h2", "h3"])],
                },
            )
        ]


class Image(File):
    """PNG/JPEG/GIF/BMP/TIFF/WebP image."""

    def parse(self, file_bytes: bytes, file_name: str, mime_type: str) -> list[Section]:
        """Decode an image and extract its metadata (and optional OCR text).

        Args:
            file_bytes: Raw image bytes.
            file_name: Original filename (unused).
            mime_type: MIME type (unused).

        Returns:
            A single-element list with one :class:`Section`.
            ``source_location`` is ``"image"``. ``metadata`` carries
            ``format``, ``size``, ``mode``, and ``exif`` (a dict of
            EXIF tag → stringified value, empty when EXIF is absent).
            ``text`` contains the OCR result, or ``""`` if
            :mod:`pytesseract` is unavailable or fails.
        """
        image = PillowImage.open(BytesIO(file_bytes))
        pytesseract_module, import_error = capture(import_module, "pytesseract")
        text = ""
        if import_error is None:
            ocr_text, ocr_error = capture(pytesseract_module.image_to_string, image)
            if ocr_error is None:
                text = ocr_text
        exif_data = image.getexif() if hasattr(image, "getexif") else {}
        metadata = {
            "format": image.format,
            "size": image.size,
            "mode": image.mode,
            "exif": {k: str(v) for k, v in exif_data.items()} if exif_data else {},
        }
        return [
            Section(
                section_index=0,
                source_location="image",
                text=text,
                metadata=metadata,
            )
        ]


class Office(File):
    """DOCX/XLSX/PPTX (also DOC/XLS/PPT) document."""

    def parse(self, file_bytes: bytes, file_name: str, mime_type: str) -> list[Section]:
        """Dispatch an Office document to its format-specific parser.

        Args:
            file_bytes: Raw file bytes.
            file_name: Used as a fallback when ``mime_type`` is empty
                or unrecognised (extension-derived dispatch).
            mime_type: MIME type used to choose the parser.

        Returns:
            A list of :class:`Section`:

            * DOCX: one section containing every paragraph joined
              with newlines; ``source_location="document"``.
            * XLSX: one section per worksheet, formatted as
              ``" | "``-joined rows; ``source_location="worksheet NAME"``.
            * PPTX: one section per slide; ``source_location="slide N"``.

            Empty when the extension/MIME pair is not an Office type.
        """
        ext = file_name.lower().rsplit(".", 1)[-1] if "." in file_name else ""
        sections: list[Section] = []

        if mime_type in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ) or ext in ("docx", "doc"):
            doc = Document(io.BytesIO(file_bytes))
            text_parts = [para.text for para in doc.paragraphs]
            sections.append(
                Section(
                    section_index=0,
                    source_location="document",
                    text="\n".join(text_parts),
                    metadata={"tables": len(doc.tables)},
                )
            )

        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ) or ext in ("xlsx", "xls"):
            wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            for i, ws_name in enumerate(wb.sheetnames, start=1):
                ws = wb[ws_name]
                rows: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) if c is not None else "" for c in row)
                    rows.append(row_text)
                sections.append(
                    Section(
                        section_index=i,
                        source_location=f"worksheet {ws_name}",
                        text="\n".join(rows),
                        metadata={"sheet_name": ws_name},
                    )
                )
            wb.close()

        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ) or ext in ("pptx", "ppt"):
            prs = Presentation(io.BytesIO(file_bytes))
            for i, slide in enumerate(prs.slides, start=1):
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                sections.append(
                    Section(
                        section_index=i,
                        source_location=f"slide {i}",
                        text="\n".join(texts),
                        metadata={"slide_number": i},
                    )
                )

        return sections


class Csv(File):
    """CSV file (UTF-8 decoded, no structural splitting)."""

    def parse(self, file_bytes: bytes, file_name: str, mime_type: str) -> list[Section]:
        """Decode CSV bytes as UTF-8 and emit one whole-file section.

        Args:
            file_bytes: Raw CSV bytes.
            file_name: Original filename (unused).
            mime_type: MIME type (unused).

        Returns:
            A single-element list containing the full decoded text.
            Invalid byte sequences are replaced with the Unicode
            replacement character (``errors="replace"``).
        """
        text = file_bytes.decode("utf-8", errors="replace")
        return [
            Section(
                section_index=0,
                source_location="full file",
                text=text,
                metadata={},
            )
        ]


class Txt(File):
    """Plain text file."""

    def parse(self, file_bytes: bytes, file_name: str, mime_type: str) -> list[Section]:
        """Decode text bytes as UTF-8 and emit one whole-file section.

        Args:
            file_bytes: Raw text bytes.
            file_name: Original filename (unused).
            mime_type: MIME type (unused).

        Returns:
            A single-element list containing the decoded text. Invalid
            byte sequences are replaced with the Unicode replacement
            character (``errors="replace"``).
        """
        text = file_bytes.decode("utf-8", errors="replace")
        return [
            Section(
                section_index=0,
                source_location="full file",
                text=text,
                metadata={},
            )
        ]


class Catalog:
    """Two-tier (MIME then extension) parser lookup with a UTF-8 fallback."""

    def __init__(self) -> None:
        """Initialise the registry and install the default parsers."""
        self.parsers: dict[str, File] = {}
        self.register_defaults()

    def register_defaults(self) -> None:
        """Install the standard parser set.

        The same :class:`File` instance is registered under
        every MIME and extension in its family (e.g. one
        :class:`Image` for png/jpeg/gif/webp/tiff). This keeps
        the registry compact and lets the same object back several
        lookups.
        """
        pdf = Pdf()
        html = HTML()
        image = Image()
        office = Office()
        csv = Csv()
        txt = Txt()

        self.register("application/pdf", pdf)
        self.register("text/html", html)
        self.register("text/plain", txt)
        self.register("text/csv", csv)
        self.register("image/png", image)
        self.register("image/jpeg", image)
        self.register("image/jpg", image)
        self.register("image/gif", image)
        self.register("image/webp", image)
        self.register("image/tiff", image)
        self.register(
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document", office
        )
        self.register("application/msword", office)
        self.register("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", office)
        self.register("application/vnd.ms-excel", office)
        self.register(
            "application/vnd.openxmlformats-officedocument.presentationml.presentation", office
        )
        self.register("application/vnd.ms-powerpoint", office)

        # Same parsers, addressed by file extension. The leading dot
        # is included so lookups are unambiguous (``".pdf"`` vs
        # ``"pdf"``).
        self.register(".pdf", pdf)
        self.register(".html", html)
        self.register(".htm", html)
        self.register(".txt", txt)
        self.register(".csv", csv)
        self.register(".png", image)
        self.register(".jpg", image)
        self.register(".jpeg", image)
        self.register(".gif", image)
        self.register(".webp", image)
        self.register(".tiff", image)
        self.register(".tif", image)
        self.register(".docx", office)
        self.register(".doc", office)
        self.register(".xlsx", office)
        self.register(".xls", office)
        self.register(".pptx", office)
        self.register(".ppt", office)

    def register(self, key: str, parser: File) -> None:
        """Register ``parser`` under ``key``.

        Args:
            key: Either a MIME type (``"text/plain"``) or an extension
                including the leading dot (``".txt"``).
            parser: The :class:`File` instance to use for that key.
        """
        self.entries[key] = parser

    def lookup(self, mime_type: str, file_name: str) -> File | None:
        """Look up a parser by MIME type then by extension.

        Args:
            mime_type: The MIME type to try first.
            file_name: Used to derive the extension fallback.

        Returns:
            The matching :class:`File`, or ``None`` if neither key is
            registered.
        """
        parser = self.entries.get(mime_type)
        if parser is not None:
            return parser
        # Extension fallback: only consult when the filename has a dot;
        # ``Path.suffix`` returns ``""`` for dot-less names.
        ext = Path(file_name).suffix.lower() if "." in file_name else ""
        return self.entries.get(ext)

    def parse(self, file_bytes: bytes, file_name: str, mime_type: str) -> list[Section]:
        """Dispatch to the appropriate parser, or fall back to UTF-8.

        Args:
            file_bytes: Raw file contents.
            file_name: Original filename.
            mime_type: MIME type reported by the validator.

        Returns:
            A list of :class:`Section`. When no parser is
            registered the function decodes the bytes as UTF-8 with
            ``errors="replace"`` and returns a single section tagged
            ``source_location="unknown"``. This silent fallback is
            intentional — it lets the pipeline gracefully accept
            unknown text-like formats — but callers that need strict
            format enforcement should validate up front.
        """
        parser = self.lookup(mime_type, file_name)
        if parser is None:
            # Forgiving fallback: decode as UTF-8 and treat the whole
            # file as one section. ``errors="replace"`` keeps the
            # decoder from raising on invalid byte sequences, at the
            # cost of substituting the Unicode replacement character.
            return [
                Section(
                    section_index=0,
                    source_location="unknown",
                    text=file_bytes.decode("utf-8", errors="replace"),
                    metadata={},
                )
            ]
        return parser.parse(file_bytes, file_name, mime_type)


# ``Catalog`` is the public class.


# The single user-facing dispatch. Callers should reach for ``parse``
# rather than instantiating their own :class:`Catalog`.
def parse(file_bytes: bytes, file_name: str, mime_type: str) -> list[Section]:
    """Parse any file via a fresh default catalog.

    This is the canonical entry point. Concrete :class:`File`
    instances and :class:`Catalog` are exposed for tests and
    custom setups; runtime ingestion should call ``parse``.

    Args:
        file_bytes: Raw file contents.
        file_name: Original filename.
        mime_type: MIME type reported by the validator.

    Returns:
        A list of :class:`Section` produced by the matched parser.
    """
    return Catalog().parse(file_bytes, file_name, mime_type)

__all__ = [
    "EQUATION_BLOCK_RE",
    "FENCE_RE",
    "HEADING_RE",
    "IMAGE_RE",
    "INLINE_EQUATION_RE",
    "MAGIC_BYTES",
    "MARKER_AVAILABLE",
    "MIME_TYPES_BY_EXTENSION",
    "TABLE_LINE_RE",
    "ChunkingPlan",
    "DocumentLifecycleManager",
    "MarkdownSection",
    "MarkerConverter",
    "MarkerImportError",
    "MarkerPdfConverter",
    "PlainTextConverter",
    "build_chunk_records",
    "build_marker_converter",
    "chunk_words",
    "convert_path",
    "datetime_now_utc",
    "detect_mime_type",
    "extract_pdf_metadata",
    "extract_pdf_pages",
    "extract_pdf_text",
    "extract_text_from_content",
    "looks_like_pdf",
    "markdown_to_document_blocks",
    "marker_create_model_dict",
    "marker_text_from_rendered",
    "new_version",
    "normalise_markdown",
    "normalize_text",
    "select_converter_for_path",
    "self_module",
    "validate_upload",
]
