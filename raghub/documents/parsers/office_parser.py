"""Office document parser.

Routes DOCX/XLSX/PPTX (and their legacy ``.doc``/``.xls``/``.ppt``
extensions) through the appropriate library:

* DOCX/DOC — :mod:`docx`
* XLSX/XLS — :mod:`openpyxl`
* PPTX/PPT — :mod:`pptx`

Legacy binary formats rely on whichever converter is available in the
runtime environment; if the corresponding library is missing the
import will surface as a parser error.
"""

from __future__ import annotations

import io

from .base import FileParser, ParsedSection


class OfficeParser(FileParser):
    """Parser for DOCX/XLSX/PPTX (and legacy DOC/XLS/PPT) documents."""

    def parse(self, file_bytes: bytes, file_name: str, mime_type: str) -> list[ParsedSection]:
        """Dispatch an Office document to its format-specific parser.

        Args:
            file_bytes: Raw file bytes.
            file_name: Used as a fallback when ``mime_type`` is empty
                or unrecognised (extension-derived dispatch).
            mime_type: MIME type used to choose the parser.

        Returns:
            A list of :class:`ParsedSection`:

            * DOCX: one section containing every paragraph joined
              with newlines; ``source_location="document"``.
            * XLSX: one section per worksheet, formatted as
              ``" | "``-joined rows; ``source_location="worksheet NAME"``.
            * PPTX: one section per slide; ``source_location="slide N"``.

            Empty when the extension/MIME pair is not an Office type.
        """
        ext = file_name.lower().rsplit(".", 1)[-1] if "." in file_name else ""
        sections: list[ParsedSection] = []

        if mime_type in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        ) or ext in ("docx", "doc"):
            from docx import Document

            doc = Document(io.BytesIO(file_bytes))
            text_parts = [para.text for para in doc.paragraphs]
            sections.append(
                ParsedSection(
                    section_index=0,
                    source_location="document",
                    text="\n".join(text_parts),
                    metadata={"tables": len(doc.tables)},
                )
            )

        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        ) or ext in ("xlsx", "xls"):
            # ``openpyxl`` is the standard read-only XLSX driver; the
            # ``data_only=True`` flag resolves formula cells to their
            # cached values instead of returning the formula strings.
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
            for i, ws_name in enumerate(wb.sheetnames, start=1):
                ws = wb[ws_name]
                rows: list[str] = []
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(c) if c is not None else "" for c in row)
                    rows.append(row_text)
                sections.append(
                    ParsedSection(
                        section_index=i,
                        source_location=f"worksheet {ws_name}",
                        text="\n".join(rows),
                        metadata={"sheet_name": ws_name},
                    )
                )
            wb.close()

        elif mime_type in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        ) or ext in ("pptx", "ppt"):
            from pptx import Presentation

            prs = Presentation(io.BytesIO(file_bytes))
            for i, slide in enumerate(prs.slides, start=1):
                texts = [shape.text for shape in slide.shapes if hasattr(shape, "text")]
                sections.append(
                    ParsedSection(
                        section_index=i,
                        source_location=f"slide {i}",
                        text="\n".join(texts),
                        metadata={"slide_number": i},
                    )
                )

        return sections
