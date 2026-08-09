"""Document parsers and converter catalog.

The parser classes (Markdown, plain text, Marker, etc.) and the
:class:`Catalog` factory live here. Validation, chunking, and
lifecycle helpers live in :mod:`raghub.lifecycle`.
"""

from __future__ import annotations

import io
from abc import ABC, abstractmethod
from dataclasses import dataclass
from importlib import import_module
from io import BytesIO
from pathlib import Path
from typing import Any

from raghub.runtime import capture
from raghub.errors import MissingDepError
from raghub.lifecycle import (
    ChunkingPlan,
    Lifecycle,
    Marker,
    chunk_words,
    normalize_text,
    validate_upload,
)

__all__ = [
    "HTML",
    "Catalog",
    "ChunkingPlan",
    "Csv",
    "Image",
    "Lifecycle",
    "Marker",
    "Office",
    "ParsedSection",
    "Pdf",
    "Txt",
    "chunk_words",
    "normalize_text",
    "parse",
    "validate_upload",
]


@dataclass(frozen=True)
class ParsedSection:
    """A single parsed chunk of a document.

    Attributes:
        section_index: 0-based ordinal of this section within the file
            (e.g. PDF page number minus 1, or 0 for whole-file formats).
        source_location: Human-readable location string used as the
            ``source_location`` field on chunk records.
        text: The extracted text content.
        metadata: Format-specific metadata (e.g. PDF ``/Title``,
            ``/Author``). Optional; defaults to an empty dict.

    """

    section_index: int
    source_location: str
    text: str
    metadata: dict[str, Any]


class File(ABC):
    """Abstract base for all document parsers."""

    @abstractmethod
    def parse(self, file_bytes: bytes, file_name: str, mime_type: str) -> list[ParsedSection]:
        """Parse ``file_bytes`` into a list of :class:`Section`.

        Args:
            file_bytes: Raw file contents.
            file_name: Original filename; useful for format-specific
                hints (e.g. extension-based fallbacks).
            mime_type: The MIME type reported by the validator.

        Returns:
            A list of :class:`Section` objects. Empty when the
            parser finds no extractable text.

        """


class Pdf(File):
    """PDF file using :mod:`pypdf`."""

    @staticmethod
    def parse(file_bytes: bytes, file_name: str, mime_type: str) -> list[ParsedSection]:
        """Parse a PDF into one section per page.

        Args:
            file_bytes: Raw PDF bytes.
            file_name: Original filename (unused beyond diagnostics).
            mime_type: MIME type (unused; the parser always uses
                :class:`pypdf.PdfReader`).

        Returns:
            A list of :class:`Section`, one per page. Empty
            strings are returned for image-only pages rather than
            raising. The section's ``section_index`` is the 1-based
            page number and ``source_location`` is ``"page N"``. The
            metadata dict contains ``width`` and ``height`` (from the
            page's media box) when available.

        """
        try:
            from pypdf import PdfReader
        except ImportError:
            raise MissingDepError(
                "pypdf",
                "pip install raghub[pdf]",
            ) from None
        reader = PdfReader(BytesIO(file_bytes))
        sections: list[ParsedSection] = []
        for i, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            sections.append(
                ParsedSection(
                    section_index=i,
                    source_location=f"page {i}",
                    text=text,
                    metadata={
                        "width": page.mediabox.width,
                        "height": page.mediabox.height,
                    }
                    if page.mediabox
                    else {},
                )
            )
        return sections


class HTML(File):
    """HTML file using :mod:`BeautifulSoup`."""

    @staticmethod
    def parse(file_bytes: bytes, file_name: str, mime_type: str) -> list[ParsedSection]:
        """Parse an HTML document into a single section.

        Args:
            file_bytes: Raw HTML bytes.
            file_name: Original filename (unused).
            mime_type: MIME type (unused).

        Returns:
            A single-element list with one :class:`Section`
            containing the concatenated body text. ``section_index``
            is 0 and ``source_location`` is ``"full file"``. The
            section metadata includes a ``headings`` list with the
            text of every ``<h1>``, ``<h2>``, and ``<h3>`` element.

        """
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise MissingDepError(
                "beautifulsoup4",
                "pip install raghub[docs]",
            ) from None
        soup = BeautifulSoup(file_bytes, "lxml")
        body = soup.find("body") or soup
        text = body.get_text(separator=" ", strip=True)
        return [
            ParsedSection(
                section_index=0,
                source_location="full file",
                text=text,
                metadata={
                    "headings": [h.get_text(strip=True) for h in soup.find_all(["h1", "h2", "h3"])],
                },
            )
        ]


class Image(File):
    """PNG/JPEG/GIF/BMP/TIFF/WebP image."""

    @staticmethod
    def parse(file_bytes: bytes, file_name: str, mime_type: str) -> list[ParsedSection]:
        """Decode an image and extract its metadata (and optional OCR text).

        Args:
            file_bytes: Raw image bytes.
            file_name: Original filename (unused).
            mime_type: MIME type (unused).

        Returns:
            A single-element list with one :class:`Section`.
            ``source_location`` is ``"image"``. ``metadata`` carries
            ``format``, ``size``, ``mode``, and ``exif`` (a dict of
            EXIF tag → stringified value, empty when EXIF is absent).
            ``text`` contains the OCR result, or ``""`` if
            :mod:`pytesseract` is unavailable or fails.

        """
        try:
            from PIL import Image as PillowImage
        except ImportError:
            raise MissingDepError(
                "Pillow",
                "pip install raghub[docs]",
            ) from None
        image = PillowImage.open(BytesIO(file_bytes))
        pytesseract_module, import_error = capture(import_module, "pytesseract")
        text = ""
        if import_error is None:
            ocr_text, ocr_error = capture(pytesseract_module.image_to_string, image)
            if ocr_error is None:
                text = ocr_text
        exif_data = image.getexif() if hasattr(image, "getexif") else {}
        metadata = {
            "format": image.format,
            "size": image.size,
            "mode": image.mode,
            "exif": {k: str(v) for k, v in exif_data.items()} if exif_data else {},
        }
        return [
            ParsedSection(
                section_index=0,
                source_location="image",
                text=text,
                metadata=metadata,
            )
        ]


class Office(File):
    """DOCX/XLSX/PPTX (also DOC/XLS/PPT) document."""

    @staticmethod
    def parse(file_bytes: bytes, file_name: str, mime_type: str) -> list[ParsedSection]:
        """Dispatch an Office document to its format-specific parser.

        Returns one ``ParsedSection`` per format-specific unit (DOCX
        document, XLSX worksheet, PPTX slide). Returns an empty list
        when the extension / MIME pair is not an Office type.
        """
        ext = file_name.lower().rsplit(".", 1)[-1] if "." in file_name else ""
        sections: list[ParsedSection] = []

        if Office.is_docx(mime_type, ext):
            sections.append(Office.parse_docx(file_bytes))
        elif Office.is_xlsx(mime_type, ext):
            sections.extend(Office.parse_xlsx(file_bytes))
        elif Office.is_pptx(mime_type, ext):
            sections.extend(Office.parse_pptx(file_bytes))

        return sections

    @staticmethod
    def is_docx(mime_type: str, ext: str) -> bool:
        return mime_type in {
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword",
        } or ext in {"docx", "doc"}

    @staticmethod
    def is_xlsx(mime_type: str, ext: str) -> bool:
        return mime_type in {
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel",
        } or ext in {"xlsx", "xls"}

    @staticmethod
    def is_pptx(mime_type: str, ext: str) -> bool:
        return mime_type in {
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint",
        } or ext in {"pptx", "ppt"}

    @staticmethod
    def parse_docx(file_bytes: bytes) -> ParsedSection:
        try:
            from docx import Document
        except ImportError as exc:
            raise MissingDepError(
                "python-docx", "pip install raghub[docs]"
            ) from exc
        doc = Document(io.BytesIO(file_bytes))
        text_parts = [para.text for para in doc.paragraphs]
        return ParsedSection(
            section_index=0,
            source_location="document",
            text="\n".join(text_parts),
            metadata={"tables": len(doc.tables)},
        )

    @staticmethod
    def parse_xlsx(file_bytes: bytes) -> list[ParsedSection]:
        try:
            from openpyxl import load_workbook
        except ImportError as exc:
            raise MissingDepError(
                "openpyxl", "pip install raghub[docs]"
            ) from exc
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
        result: list[ParsedSection] = []
        for i, ws_name in enumerate(wb.sheetnames, start=1):
            ws = wb[ws_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(
                    str(c) if c is not None else "" for c in row
                )
                rows.append(row_text)
            result.append(
                ParsedSection(
                    section_index=i,
                    source_location=f"worksheet {ws_name}",
                    text="\n".join(rows),
                    metadata={"sheet_name": ws_name},
                )
            )
        wb.close()
        return result

    @staticmethod
    def parse_pptx(file_bytes: bytes) -> list[ParsedSection]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise MissingDepError(
                "python-pptx", "pip install raghub[docs]"
            ) from exc
        prs = Presentation(io.BytesIO(file_bytes))
        result: list[ParsedSection] = []
        for i, slide in enumerate(prs.slides, start=1):
            texts = [
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text")
            ]
            result.append(
                ParsedSection(
                    section_index=i,
                    source_location=f"slide {i}",
                    text="\n".join(texts),
                    metadata={"slide_number": i},
                )
            )
        return result


class Csv(File):
    """CSV file (UTF-8 decoded, no structural splitting)."""

    @staticmethod
    def parse(file_bytes: bytes, file_name: str, mime_type: str) -> list[ParsedSection]:
        """Decode CSV bytes as UTF-8 and emit one whole-file section.

        Args:
            file_bytes: Raw CSV bytes.
            file_name: Original filename (unused).
            mime_type: MIME type (unused).

        Returns:
            A single-element list containing the full decoded text.
            Invalid byte sequences are replaced with the Unicode
            replacement character (``errors="replace"``).

        """
        text = file_bytes.decode("utf-8", errors="replace")
        return [
            ParsedSection(
                section_index=0,
                source_location="full file",
                text=text,
                metadata={},
            )
        ]


class Txt(File):
    """Plain text file."""

    @staticmethod
    def parse(file_bytes: bytes, file_name: str, mime_type: str) -> list[ParsedSection]:
        """Decode text bytes as UTF-8 and emit one whole-file section.

        Args:
            file_bytes: Raw text bytes.
            file_name: Original filename (unused).
            mime_type: MIME type (unused).

        Returns:
            A single-element list containing the decoded text. Invalid
            byte sequences are replaced with the Unicode replacement
            character (``errors="replace"``).

        """
        text = file_bytes.decode("utf-8", errors="replace")
        return [
            ParsedSection(
                section_index=0,
                source_location="full file",
                text=text,
                metadata={},
            )
        ]


class Catalog:
    """Two-tier (MIME then extension) parser lookup with a UTF-8 fallback."""

    def __init__(self) -> None:
        """Initialise the registry and install the default parsers."""
        self.entries: dict[str, File] = {}
        self.register_defaults()

    def register_defaults(self) -> None:
        """Install the standard parser set.

        The same :class:`File` instance is registered under
        every MIME and extension in its family (e.g. one
        :class:`Image` for png/jpeg/gif/webp/tiff). This keeps
        the registry compact and lets the same object back several
        lookups.
        """
        parsers = self._build_default_parsers()
        for mime, parser in parsers:
            self.register(mime, parser)
        for ext, parser in self._build_default_extensions(parsers):
            self.register(ext, parser)

    @staticmethod
    def _build_default_parsers() -> list[tuple[str, File]]:
        """Return ``(mime, File)`` pairs for every built-in parser type."""
        return [
            ("application/pdf", Pdf()),
            ("text/html", HTML()),
            ("text/plain", Txt()),
            ("text/csv", Csv()),
            ("image/png", Image()),
            ("image/jpeg", Image()),
            ("image/jpg", Image()),
            ("image/gif", Image()),
            ("image/webp", Image()),
            ("image/tiff", Image()),
            (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                Office(),
            ),
            ("application/msword", Office()),
            (
                "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                Office(),
            ),
            ("application/vnd.ms-excel", Office()),
            (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                Office(),
            ),
            ("application/vnd.ms-powerpoint", Office()),
        ]

    @staticmethod
    def _build_default_extensions(
        parsers: list[tuple[str, File]]
    ) -> list[tuple[str, File]]:
        """Return ``(extension, File)`` pairs (with leading dot) for built-in parsers."""
        extension_map: dict[str, File] = {
            ".pdf": parsers[0][1],
            ".html": parsers[1][1],
            ".htm": parsers[1][1],
            ".txt": parsers[2][1],
            ".csv": parsers[3][1],
            ".png": parsers[4][1],
            ".jpg": parsers[5][1],
            ".jpeg": parsers[6][1],
            ".gif": parsers[7][1],
            ".webp": parsers[8][1],
            ".tiff": parsers[9][1],
            ".tif": parsers[9][1],
            ".docx": parsers[10][1],
            ".doc": parsers[11][1],
            ".xlsx": parsers[12][1],
            ".xls": parsers[13][1],
            ".pptx": parsers[14][1],
            ".ppt": parsers[15][1],
        }
        return list(extension_map.items())

    def register(self, key: str, parser: File) -> None:
        """Register ``parser`` under ``key``.

        Args:
            key: Either a MIME type (``"text/plain"``) or an extension
                including the leading dot (``".txt"``).
            parser: The :class:`File` instance to use for that key.

        """
        self.entries[key] = parser

    def lookup(self, mime_type: str, file_name: str) -> File | None:
        """Look up a parser by MIME type then by extension.

        Args:
            mime_type: The MIME type to try first.
            file_name: Used to derive the extension fallback.

        Returns:
            The matching :class:`File`, or ``None`` if neither key is
            registered.

        """
        parser = self.entries.get(mime_type)
        if parser is not None:
            return parser
        # Extension fallback: only consult when the filename has a dot;
        # ``Path.suffix`` returns ``""`` for dot-less names.
        ext = Path(file_name).suffix.lower() if "." in file_name else ""
        return self.entries.get(ext)

    def parse(self, file_bytes: bytes, file_name: str, mime_type: str) -> list[ParsedSection]:
        """Dispatch to the appropriate parser, or fall back to UTF-8.

        Args:
            file_bytes: Raw file contents.
            file_name: Original filename.
            mime_type: MIME type reported by the validator.

        Returns:
            A list of :class:`Section`. When no parser is
            registered the function decodes the bytes as UTF-8 with
            ``errors="replace"`` and returns a single section tagged
            ``source_location="unknown"``. This silent fallback is
            intentional — it lets the pipeline gracefully accept
            unknown text-like formats — but callers that need strict
            format enforcement should validate up front.

        """
        parser = self.lookup(mime_type, file_name)
        if parser is None:
            # Forgiving fallback: decode as UTF-8 and treat the whole
            # file as one section. ``errors="replace"`` keeps the
            # decoder from raising on invalid byte sequences, at the
            # cost of substituting the Unicode replacement character.
            return [
                ParsedSection(
                    section_index=0,
                    source_location="unknown",
                    text=file_bytes.decode("utf-8", errors="replace"),
                    metadata={},
                )
            ]
        return parser.parse(file_bytes, file_name, mime_type)


# ``Catalog`` is the public class.


# The single user-facing dispatch. Callers should reach for ``parse``
# rather than instantiating their own :class:`Catalog`.
def parse(file_bytes: bytes, file_name: str, mime_type: str) -> list[ParsedSection]:
    """Parse any file via a fresh default catalog.

    This is the canonical entry point. Concrete :class:`File`
    instances and :class:`Catalog` are exposed for tests and
    custom setups; runtime ingestion should call ``parse``.

    Args:
        file_bytes: Raw file contents.
        file_name: Original filename.
        mime_type: MIME type reported by the validator.

    Returns:
        A list of :class:`Section` produced by the matched parser.

    """
    return Catalog().parse(file_bytes, file_name, mime_type)
