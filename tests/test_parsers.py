from __future__ import annotations

from pathlib import Path

import pytest

from raghub.documents.parsers import ParserRegistry
from raghub.documents.parsers.base import ParsedSection
from raghub.documents.parsers.pdf_parser import PdfParser
from raghub.documents.parsers.html_parser import HtmlParser
from raghub.documents.parsers.image_parser import ImageParser
from raghub.documents.parsers.office_parser import OfficeParser
from raghub.documents.parsers.csv_parser import CsvParser
from raghub.documents.parsers.txt_parser import TxtParser


class TestPdfParser:
    def test_extracts_text_by_page(self):
        parser = PdfParser()
        pdf_bytes = _make_minimal_pdf()
        sections = parser.parse(pdf_bytes, "test.pdf", "application/pdf")
        assert len(sections) >= 1
        assert sections[0].source_location.startswith("page")
        assert sections[0].section_index >= 1


class TestHtmlParser:
    def test_extracts_text(self):
        parser = HtmlParser()
        html = b"<html><body><p>Hello world</p></body></html>"
        sections = parser.parse(html, "test.html", "text/html")
        assert len(sections) == 1
        assert "Hello" in sections[0].text
        assert sections[0].source_location == "full file"

    def test_extracts_headings_in_metadata(self):
        parser = HtmlParser()
        html = b"<html><body><h1>Title</h1><h2>Sub</h2><p>Text</p></body></html>"
        sections = parser.parse(html, "test.html", "text/html")
        assert "headings" in sections[0].metadata
        assert "Title" in sections[0].metadata["headings"]


class TestImageParser:
    def test_returns_exif_metadata_for_png(self):
        parser = ImageParser()
        img_bytes = _make_minimal_png()
        sections = parser.parse(img_bytes, "test.png", "image/png")
        assert len(sections) == 1
        assert sections[0].source_location == "image"
        assert "format" in sections[0].metadata
        assert sections[0].metadata["format"] == "PNG"


class TestCsvParser:
    def test_returns_full_content(self):
        parser = CsvParser()
        csv = b"name,age\nAlice,30\nBob,25\n"
        sections = parser.parse(csv, "test.csv", "text/csv")
        assert len(sections) == 1
        assert "Alice" in sections[0].text
        assert sections[0].source_location == "full file"


class TestTxtParser:
    def test_returns_text(self):
        parser = TxtParser()
        content = b"Hello, world!"
        sections = parser.parse(content, "test.txt", "text/plain")
        assert len(sections) == 1
        assert sections[0].text == "Hello, world!"
        assert sections[0].source_location == "full file"


class TestOfficeParser:
    def test_docx_returns_sections(self):
        parser = OfficeParser()
        docx_bytes = _make_minimal_docx()
        sections = parser.parse(docx_bytes, "test.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document")
        assert len(sections) >= 1
        assert sections[0].source_location == "document"

    def test_xlsx_returns_worksheets(self):
        parser = OfficeParser()
        xlsx_bytes = _make_minimal_xlsx()
        sections = parser.parse(xlsx_bytes, "test.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
        assert len(sections) >= 1
        assert "worksheet" in sections[0].source_location

    def test_pptx_returns_slides(self):
        parser = OfficeParser()
        pptx_bytes = _make_minimal_pptx()
        sections = parser.parse(pptx_bytes, "test.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation")
        assert len(sections) >= 1
        assert "slide" in sections[0].source_location


class TestParserRegistry:
    def test_registry_returns_parser_by_mime(self):
        registry = ParserRegistry()
        parser = registry.get_parser("application/pdf", "doc.pdf")
        assert parser is not None

    def test_registry_returns_parser_by_extension(self):
        registry = ParserRegistry()
        parser = registry.get_parser("unknown/type", "doc.pdf")
        assert parser is not None

    def test_registry_returns_none_for_unknown(self):
        registry = ParserRegistry()
        parser = registry.get_parser("unknown/type", "doc.xyz")
        assert parser is None

    def test_parse_dispatches_correctly(self):
        registry = ParserRegistry()
        csv = b"a,b\n1,2\n"
        sections = registry.parse(csv, "data.csv", "text/csv")
        assert len(sections) == 1
        assert "1" in sections[0].text


def _make_minimal_pdf() -> bytes:
    from io import BytesIO
    from pypdf import PdfWriter
    w = PdfWriter()
    w.add_blank_page(612, 792)
    buf = BytesIO()
    w.write(buf)
    return buf.getvalue()


def _make_minimal_png() -> bytes:
    import struct
    import zlib

    def make_chunk(chunk_type: bytes, data: bytes) -> bytes:
        chunk = chunk_type + data
        return struct.pack(">I", len(data)) + chunk + struct.pack(">I", zlib.crc32(chunk) & 0xFFFFFFFF)

    signature = b"\x89PNG\r\n\x1a\n"
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr = make_chunk(b"IHDR", ihdr_data)
    raw_data = zlib.compress(b"\x00\xff\x00\x00")
    idat = make_chunk(b"IDAT", raw_data)
    iend = make_chunk(b"IEND", b"")
    return signature + ihdr + idat + iend


def _make_minimal_docx() -> bytes:
    from docx import Document
    from io import BytesIO
    doc = Document()
    doc.add_paragraph("Hello World")
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _make_minimal_xlsx() -> bytes:
    from openpyxl import Workbook
    from io import BytesIO
    wb = Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Hello"
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _make_minimal_pptx() -> bytes:
    from pptx import Presentation
    from io import BytesIO
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Hello PowerPoint"
    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
